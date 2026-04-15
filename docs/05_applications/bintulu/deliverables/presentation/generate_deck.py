#!/usr/bin/env python3
"""Generate Bintulu Smart City presentation deck.

Produces a ~20-slide professional 16:9 PPTX for Bintulu Development Authority
(BDA), Malaysia. VIETSOL (with partner ESP) proposes AI-powered Smart Parking
and AI Traffic Light solutions.

Output: bintulu_smart_city_deck.pptx (same directory as this script).
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colour palette
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
MEDIUM_BLUE = RGBColor(0x2E, 0x50, 0x90)
ACCENT_BLUE = RGBColor(0x44, 0x72, 0xC4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GREEN = RGBColor(0x2D, 0x8B, 0x4E)
RED = RGBColor(0xC0, 0x39, 0x2B)
GOLD = RGBColor(0xD4, 0xA0, 0x1E)

FONT_FAMILY = "Calibri"
FOOTER_TEXT = "Strictly Confidential | VIETSOL + ESP | March 2026"


# ---------------------------------------------------------------------------
# Helper: font styling
# ---------------------------------------------------------------------------
def _style_run(run, *, size=Pt(18), bold=False, italic=False, color=DARK_GRAY,
               font_name=FONT_FAMILY):
    """Apply common font styling to a run."""
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name


def _style_paragraph(para, *, alignment=PP_ALIGN.LEFT, space_after=Pt(6),
                     space_before=Pt(0), level=0):
    """Apply paragraph-level formatting."""
    para.alignment = alignment
    para.space_after = space_after
    para.space_before = space_before
    para.level = level


# ---------------------------------------------------------------------------
# Helper: add_footer
# ---------------------------------------------------------------------------
def add_footer(slide, text=FOOTER_TEXT):
    """Add a footer text box at the bottom of *slide*."""
    left = Inches(0.5)
    top = SLIDE_HEIGHT - Inches(0.45)
    width = SLIDE_WIDTH - Inches(1.0)
    height = Inches(0.35)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    _style_run(run, size=Pt(9), color=RGBColor(0x99, 0x99, 0x99), italic=True)


# ---------------------------------------------------------------------------
# Helper: dark-background rectangle (full slide)
# ---------------------------------------------------------------------------
def _add_bg_rect(slide, color=DARK_BLUE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Send to back by moving to index 0 in the shape tree
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return shape


# ---------------------------------------------------------------------------
# Helper: add accent bar under title
# ---------------------------------------------------------------------------
def _add_accent_bar(slide, top, color=ACCENT_BLUE, width=Inches(2.0)):
    left = Inches(0.75)
    height = Inches(0.06)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def add_title_slide(prs, title, subtitle, extra_line=None):
    """Slide with centred title + subtitle on dark-blue background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg_rect(slide, DARK_BLUE)

    # Title
    left = Inches(1.0)
    top = Inches(2.0)
    width = SLIDE_WIDTH - Inches(2.0)
    height = Inches(1.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    _style_run(run, size=Pt(42), bold=True, color=WHITE)

    # Subtitle
    top2 = Inches(4.0)
    height2 = Inches(1.2)
    txBox2 = slide.shapes.add_textbox(left, top2, width, height2)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = subtitle
    _style_run(run2, size=Pt(22), color=RGBColor(0xBB, 0xCC, 0xDD))

    if extra_line:
        p3 = tf2.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(12)
        run3 = p3.add_run()
        run3.text = extra_line
        _style_run(run3, size=Pt(20), bold=True, color=ACCENT_BLUE)

    # Accent line
    bar_w = Inches(3)
    bar_left = (SLIDE_WIDTH - bar_w) // 2
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, bar_left, Inches(3.85), bar_w, Inches(0.04),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

    return slide


def add_content_slide(prs, title, bullet_points, *, footer=True):
    """Slide with a title bar and bulleted content."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title bar background
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, Inches(1.15),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.18), SLIDE_WIDTH - Inches(1.5), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    _style_run(run, size=Pt(32), bold=True, color=WHITE)

    # Accent bar
    _add_accent_bar(slide, Inches(1.15), ACCENT_BLUE, Inches(2.5))

    # Body
    body_left = Inches(0.9)
    body_top = Inches(1.55)
    body_w = SLIDE_WIDTH - Inches(1.8)
    body_h = SLIDE_HEIGHT - Inches(2.3)
    txBox2 = slide.shapes.add_textbox(body_left, body_top, body_w, body_h)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, item in enumerate(bullet_points):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        _style_paragraph(p, space_after=Pt(10), space_before=Pt(4))

        # Support bold headers via "**Header** rest" syntax
        if "**" in item:
            parts = item.split("**")
            for j, part in enumerate(parts):
                if not part:
                    continue
                run = p.add_run()
                if j % 2 == 1:  # odd index = bold part
                    run.text = part
                    _style_run(run, size=Pt(20), bold=True, color=DARK_BLUE)
                else:
                    run.text = part
                    _style_run(run, size=Pt(20), color=DARK_GRAY)
        else:
            run = p.add_run()
            run.text = item
            _style_run(run, size=Pt(20), color=DARK_GRAY)

    if footer:
        add_footer(slide)
    return slide


def add_table_slide(prs, title, headers, rows, *, footer=True,
                    col_widths=None, highlight_rows=None, bottom_note=None,
                    font_size=Pt(13), header_font_size=Pt(14)):
    """Slide with title + styled table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, Inches(1.15),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.18), SLIDE_WIDTH - Inches(1.5), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _style_run(run, size=Pt(32), bold=True, color=WHITE)

    _add_accent_bar(slide, Inches(1.15), ACCENT_BLUE, Inches(2.5))

    # Table dimensions
    num_rows = len(rows) + 1
    num_cols = len(headers)
    tbl_left = Inches(0.6)
    tbl_top = Inches(1.5)
    tbl_width = SLIDE_WIDTH - Inches(1.2)
    tbl_height = min(Inches(0.45) * num_rows, SLIDE_HEIGHT - Inches(2.5))

    table_shape = slide.shapes.add_table(num_rows, num_cols, tbl_left, tbl_top, tbl_width, tbl_height)
    table = table_shape.table

    # Column widths
    if col_widths:
        for idx, w in enumerate(col_widths):
            table.columns[idx].width = w
    else:
        default_w = int(tbl_width / num_cols)
        for idx in range(num_cols):
            table.columns[idx].width = default_w

    highlight_rows = highlight_rows or []

    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = MEDIUM_BLUE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        _style_run(run, size=header_font_size, bold=True, color=WHITE)

    # Data rows
    for ri, row_data in enumerate(rows):
        is_highlight = ri in highlight_rows
        for ci, val in enumerate(row_data):
            cell = table.cell(ri + 1, ci)
            if is_highlight:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xE8)
            elif ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            cell.text_frame.word_wrap = True

            # Bold detection
            is_bold = val.startswith("**") and val.endswith("**")
            clean = val.strip("*") if is_bold else val
            run = p.add_run()
            run.text = clean
            clr = GREEN if is_bold and is_highlight else DARK_GRAY
            _style_run(run, size=font_size, bold=is_bold, color=clr)

    # Bottom note
    if bottom_note:
        note_top = tbl_top + tbl_height + Inches(0.15)
        txB = slide.shapes.add_textbox(tbl_left, note_top, tbl_width, Inches(0.5))
        tf2 = txB.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = bottom_note
        _style_run(r2, size=Pt(16), bold=True, color=GREEN)

    if footer:
        add_footer(slide)
    return slide


def add_two_column_slide(prs, title, left_items, right_items, *,
                         left_title="", right_title="", footer=True):
    """Two-column bullet slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, Inches(1.15),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.18), SLIDE_WIDTH - Inches(1.5), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _style_run(run, size=Pt(32), bold=True, color=WHITE)

    _add_accent_bar(slide, Inches(1.15), ACCENT_BLUE, Inches(2.5))

    col_w = (SLIDE_WIDTH - Inches(2.0)) // 2
    col_h = SLIDE_HEIGHT - Inches(2.3)

    for col_idx, (items, col_title) in enumerate(
        [(left_items, left_title), (right_items, right_title)]
    ):
        left = Inches(0.9) if col_idx == 0 else Inches(0.9) + col_w + Inches(0.4)
        top = Inches(1.55)
        txB = slide.shapes.add_textbox(left, top, col_w, col_h)
        tf2 = txB.text_frame
        tf2.word_wrap = True

        if col_title:
            p0 = tf2.paragraphs[0]
            p0.alignment = PP_ALIGN.LEFT
            p0.space_after = Pt(8)
            r0 = p0.add_run()
            r0.text = col_title
            _style_run(r0, size=Pt(22), bold=True, color=MEDIUM_BLUE)
            _ = 1  # items start after title
        else:
            pass

        for i, item in enumerate(items):
            p = tf2.add_paragraph() if (i > 0 or col_title) else tf2.paragraphs[0]
            _style_paragraph(p, space_after=Pt(8))
            if "**" in item:
                parts = item.split("**")
                for j, part in enumerate(parts):
                    if not part:
                        continue
                    run = p.add_run()
                    if j % 2 == 1:
                        run.text = part
                        _style_run(run, size=Pt(18), bold=True, color=DARK_BLUE)
                    else:
                        run.text = part
                        _style_run(run, size=Pt(18), color=DARK_GRAY)
            else:
                run = p.add_run()
                run.text = item
                _style_run(run, size=Pt(18), color=DARK_GRAY)

    if footer:
        add_footer(slide)
    return slide


def _add_box(slide, left, top, width, height, text, *, fill_color=MEDIUM_BLUE,
             text_color=WHITE, font_size=Pt(14), bold=True):
    """Add a rounded rectangle with centred text. Returns the shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = DARK_BLUE
    shape.line.width = Pt(1)
    # Adjust corner rounding
    shape.adjustments[0] = 0.1

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    _style_run(run, size=font_size, bold=bold, color=text_color)
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    # Vertical centre
    tf.word_wrap = True
    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
    return shape


def _add_down_arrow_shape(slide, cx, top, length=Inches(0.5), label=None):
    """Add a downward-pointing block arrow at (cx, top)."""
    arr_w = Inches(0.35)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, cx - arr_w // 2, top, arr_w, length,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    if label:
        txB = slide.shapes.add_textbox(cx + Inches(0.25), top, Inches(2.5), Inches(0.35))
        tf = txB.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        _style_run(run, size=Pt(11), italic=True, color=MEDIUM_BLUE)


def _add_right_arrow_shape(slide, left, cy, length=Inches(0.8), label=None):
    """Add a right-pointing block arrow."""
    arr_h = Inches(0.3)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, cy - arr_h // 2, length, arr_h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    if label:
        txB = slide.shapes.add_textbox(left, cy - Inches(0.5), length, Inches(0.35))
        tf = txB.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        _style_run(run, size=Pt(10), italic=True, color=MEDIUM_BLUE)


def add_diagram_slide(prs, title, build_fn, *, footer=True):
    """Generic diagram slide — *build_fn(slide)* draws shapes on the slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, Inches(1.15),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.18), SLIDE_WIDTH - Inches(1.5), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _style_run(run, size=Pt(32), bold=True, color=WHITE)

    _add_accent_bar(slide, Inches(1.15), ACCENT_BLUE, Inches(2.5))

    build_fn(slide)

    if footer:
        add_footer(slide)
    return slide


# ===================================================================
# DIAGRAM BUILDERS (called by add_diagram_slide)
# ===================================================================

def _build_parking_architecture(slide):
    """Slide 5: three-tier parking architecture."""
    box_w = Inches(4.5)
    box_h = Inches(1.0)
    cx = (SLIDE_WIDTH - box_w) // 2

    # Row 1 — cameras
    top1 = Inches(1.8)
    _add_box(slide, cx, top1, box_w, box_h,
             "IP Cameras\n140 Occupancy + 30 ANPR",
             fill_color=MEDIUM_BLUE)

    # Arrow 1→2
    arrow1_top = top1 + box_h
    _add_down_arrow_shape(slide, cx + box_w // 2, arrow1_top,
                          Inches(0.55), label="RTSP Video Streams")

    # Row 2 — edge
    top2 = arrow1_top + Inches(0.6)
    _add_box(slide, cx, top2, box_w, box_h,
             "Edge AI Boxes\n15 Units \u2014 Local AI Processing",
             fill_color=ACCENT_BLUE)

    # Arrow 2→3
    arrow2_top = top2 + box_h
    _add_down_arrow_shape(slide, cx + box_w // 2, arrow2_top,
                          Inches(0.55), label="Metadata Only (MQTT)")

    # Row 3 — platform
    top3 = arrow2_top + Inches(0.6)
    _add_box(slide, cx, top3, box_w, box_h,
             "Central Platform\nDashboard \u2022 Analytics \u2022 Mobile App",
             fill_color=DARK_BLUE)


def _build_parking_flow(slide):
    """Slide 7: horizontal flow of how parking detection works."""
    steps = [
        "Camera\nCaptures\nFrame",
        "Edge AI\nDetects\nVehicles",
        "Bay ROI\nMatching",
        "Occupancy\nStatus\nPublished",
        "Dashboard +\nLED Signs\nUpdated",
    ]
    n = len(steps)
    box_w = Inches(2.0)
    box_h = Inches(1.3)
    gap = Inches(0.35)
    total_w = n * box_w + (n - 1) * gap
    start_left = (SLIDE_WIDTH - total_w) // 2
    top = Inches(3.2)

    colors = [MEDIUM_BLUE, ACCENT_BLUE, MEDIUM_BLUE, ACCENT_BLUE, DARK_BLUE]

    for i, (step, clr) in enumerate(zip(steps, colors)):
        left = start_left + i * (box_w + gap)
        _add_box(slide, left, top, box_w, box_h, step, fill_color=clr,
                 font_size=Pt(13))
        if i < n - 1:
            arr_left = left + box_w
            _add_right_arrow_shape(slide, arr_left, top + box_h // 2, gap)


def _build_traffic_architecture(slide):
    """Slide 12: traffic light architecture."""
    box_w = Inches(2.5)
    box_h = Inches(1.0)
    gap = Inches(0.5)
    n = 4
    labels = [
        "4 Cameras\nPer Intersection",
        "Edge AI Box\nDetection +\nClassification",
        "Signal\nController\nPhase Optimiser",
        "Traffic\nSignals\nAdaptive Timing",
    ]
    colors = [MEDIUM_BLUE, ACCENT_BLUE, MEDIUM_BLUE, DARK_BLUE]
    total_w = n * box_w + (n - 1) * gap
    start_left = (SLIDE_WIDTH - total_w) // 2
    top = Inches(2.5)

    for i, (lab, clr) in enumerate(zip(labels, colors)):
        left = start_left + i * (box_w + gap)
        _add_box(slide, left, top, box_w, box_h, lab, fill_color=clr,
                 font_size=Pt(12))
        if i < n - 1:
            arr_left = left + box_w
            _add_right_arrow_shape(slide, arr_left, top + box_h // 2, gap)

    # Arrow labels
    arrow_labels = ["RTSP", "Queue Data", "NTCIP"]
    for i, alab in enumerate(arrow_labels):
        left = start_left + i * (box_w + gap) + box_w
        txB = slide.shapes.add_textbox(left, top - Inches(0.35), gap, Inches(0.3))
        tf = txB.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = alab
        _style_run(run, size=Pt(10), italic=True, color=MEDIUM_BLUE)

    # Failsafe note
    note_top = top + box_h + Inches(0.7)
    txB2 = slide.shapes.add_textbox(Inches(2.0), note_top, SLIDE_WIDTH - Inches(4.0), Inches(0.6))
    tf2 = txB2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "\u26a0  Failsafe: automatic revert to fixed timing on any AI component failure"
    _style_run(r2, size=Pt(14), bold=True, color=RED)


def _build_shared_platform(slide):
    """Slide 16: shared platform diagram."""
    # Left column — sources
    src_w = Inches(3.0)
    src_h = Inches(0.9)
    src_left = Inches(0.6)
    src_top1 = Inches(2.0)
    src_top2 = Inches(3.4)

    _add_box(slide, src_left, src_top1, src_w, src_h,
             "Parking Zones\n7,000 Bays", fill_color=MEDIUM_BLUE, font_size=Pt(13))
    _add_box(slide, src_left, src_top2, src_w, src_h,
             "Intersections\nTraffic Signals", fill_color=MEDIUM_BLUE, font_size=Pt(13))

    # Middle — edge
    mid_left = Inches(4.5)
    mid_top = Inches(2.5)
    mid_w = Inches(2.5)
    mid_h = Inches(1.3)
    _add_box(slide, mid_left, mid_top, mid_w, mid_h,
             "Edge AI Boxes\nLocal Processing", fill_color=ACCENT_BLUE, font_size=Pt(13))

    # Arrows left→mid
    _add_right_arrow_shape(slide, src_left + src_w, src_top1 + src_h // 2, Inches(0.6))
    _add_right_arrow_shape(slide, src_left + src_w, src_top2 + src_h // 2, Inches(0.6))

    # Right — platform
    rt_left = Inches(7.9)
    rt_top = Inches(2.0)
    rt_w = Inches(3.0)
    rt_h = Inches(1.0)
    _add_box(slide, rt_left, rt_top, rt_w, rt_h,
             "Central Management\nPlatform", fill_color=DARK_BLUE, font_size=Pt(13))

    _add_right_arrow_shape(slide, mid_left + mid_w, mid_top + mid_h // 2, Inches(0.6))

    # Outputs
    outputs = ["Dashboard", "Mobile App", "LED Signs", "Signal Control"]
    out_w = Inches(2.2)
    out_h = Inches(0.55)
    out_top = Inches(3.5)
    out_gap = Inches(0.15)
    for i, out in enumerate(outputs):
        t = out_top + i * (out_h + out_gap)
        _add_box(slide, rt_left + Inches(0.4), t, out_w, out_h, out,
                 fill_color=LIGHT_GRAY, text_color=DARK_BLUE, font_size=Pt(12))


# ===================================================================
# MAIN — build all 20 slides
# ===================================================================
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ------------------------------------------------------------------
    # SLIDE 1 — Title
    # ------------------------------------------------------------------
    add_title_slide(
        prs,
        "Bintulu Smart City\nAI-Powered Urban Solutions",
        "Prepared for Bintulu Development Authority (BDA) | March 2026",
        extra_line="VIETSOL + ESP",
    )

    # ------------------------------------------------------------------
    # SLIDE 2 — Agenda
    # ------------------------------------------------------------------
    add_content_slide(prs, "Today's Presentation", [
        "\u2776  Smart Parking Solution (UC2) \u2014 7,000 bays, AI-powered occupancy & ANPR",
        "\u2777  AI Traffic Light Solution (UC1) \u2014 adaptive signal control",
        "\u2778  Shared Platform Architecture \u2014 unified management across use cases",
        "\u2779  Investment & Timeline \u2014 cost breakdown, phased delivery",
        "\u277a  Next Steps \u2014 pilot selection, project kickoff",
    ])

    # ------------------------------------------------------------------
    # SLIDE 3 — Executive Summary
    # ------------------------------------------------------------------
    add_content_slide(prs, "Executive Summary", [
        "**7,000 parking bays** managed with AI-powered camera-based detection",
        "**$450K total investment** for full parking solution \u2014 3-5x cheaper than sensor-based alternatives",
        "**$20K per intersection** for AI traffic light \u2014 55% below market median",
        "**95% detection accuracy** using proven AI detection algorithms",
        "**10-month delivery** from kickoff to full operational handover",
        "**Zero AI licensing fees** \u2014 open-source models, no per-device charges",
        "**Edge processing** \u2014 all AI runs locally, only metadata sent to cloud (privacy by design)",
    ])

    # ------------------------------------------------------------------
    # SLIDE 4 — Parking Challenge
    # ------------------------------------------------------------------
    add_content_slide(prs, "The Challenge \u2014 Parking in Bintulu", [
        "**7,000 parking bays** across multiple zones with no real-time visibility",
        "**Manual enforcement** \u2014 labour-intensive patrols, inconsistent coverage",
        "**No occupancy data** \u2014 drivers circle looking for bays, increasing congestion",
        "**Revenue leakage** \u2014 MYR 0.50/hr rate with limited compliance monitoring",
        "**Safety blind spots** \u2014 no automated detection of incidents in parking areas",
        "**No analytics** \u2014 city planners lack data for capacity and demand forecasting",
    ])

    # ------------------------------------------------------------------
    # SLIDE 5 — Parking Architecture (diagram)
    # ------------------------------------------------------------------
    add_diagram_slide(prs, "Our Solution \u2014 Smart Parking Architecture",
                      _build_parking_architecture)

    # ------------------------------------------------------------------
    # SLIDE 6 — Key Features
    # ------------------------------------------------------------------
    add_content_slide(prs, "Smart Parking \u2014 Key Features", [
        "**Bay Occupancy Detection** \u2014 95% accuracy, real-time status per bay using AI vision",
        "**ANPR (Automatic Number Plate Recognition)** \u2014 95%+ recognition, entry/exit logging, violation tracking",
        "**Safety Monitoring** \u2014 loitering detection, zone intrusion alerts, fall detection",
        "**Violation Detection** \u2014 double parking, obstruction, unauthorised zone usage",
        "**Real-Time Guidance** \u2014 LED signage with available bay counts, mobile app wayfinding",
        "**Analytics Dashboard** \u2014 occupancy trends, peak-hour analysis, revenue reports, exportable data",
    ])

    # ------------------------------------------------------------------
    # SLIDE 7 — How It Works (flow diagram)
    # ------------------------------------------------------------------
    add_diagram_slide(prs, "How It Works \u2014 Detection Pipeline",
                      _build_parking_flow)

    # ------------------------------------------------------------------
    # SLIDE 8 — Cost Comparison (HERO SLIDE)
    # ------------------------------------------------------------------
    add_table_slide(
        prs,
        "Cost Comparison \u2014 Camera-Based AI vs. Sensor-Based",
        ["Component", "Camera-Based (Ours)", "Sensor-Based"],
        [
            ["Detection Hardware", "$115K (170 cameras)", "$2.1\u2013$3.5M (7,000 sensors)"],
            ["Edge AI Compute", "$37.5K (15 AI boxes)", "$20\u2013$50K"],
            ["Installation", "$40K", "$700K\u2013$1.4M"],
            ["Software & Platform", "$180K", "$50\u2013$200K"],
            ["**Year 1 Total**", "**$450K**", "**$2.97\u2013$5.45M**"],
            ["**Per-Bay Cost**", "**~$64**", "**$400\u2013$750**"],
        ],
        highlight_rows=[4, 5],
        bottom_note="\u2b50  3\u20135x cost advantage with camera-based AI approach",
        font_size=Pt(15),
        header_font_size=Pt(16),
        col_widths=[Inches(3.5), Inches(4.5), Inches(4.5)],
    )

    # ------------------------------------------------------------------
    # SLIDE 9 — Deployment Timeline
    # ------------------------------------------------------------------
    add_table_slide(
        prs,
        "Smart Parking \u2014 Deployment Timeline (10 Months)",
        ["Phase", "Months", "Key Activities"],
        [
            ["Foundation", "1\u20134", "AI model training & validation, platform development, ANPR pipeline, hardware procurement"],
            ["Pilot", "5\u20136", "500\u20131,000 bay pilot zone, real-world testing, accuracy fine-tuning, stakeholder feedback"],
            ["Full Rollout", "7\u20138", "All 7,000 bays operational, LED signage installed, mobile app launch"],
            ["Acceptance", "9\u201310", "30-day acceptance testing, staff training, documentation, formal handover"],
        ],
        col_widths=[Inches(2.0), Inches(1.8), Inches(8.5)],
    )

    # ------------------------------------------------------------------
    # SLIDE 10 — Acceptance Criteria
    # ------------------------------------------------------------------
    add_table_slide(
        prs,
        "Smart Parking \u2014 Acceptance Criteria",
        ["Metric", "Target", "Measurement Method"],
        [
            ["Bay Occupancy Accuracy", "\u226595%", "Ground truth audit \u2014 500+ bays sampled"],
            ["ANPR Recognition (Day)", "\u226595%", "Manual plate verification \u2014 1,000+ plates"],
            ["Detection Latency", "\u22645 seconds", "Timestamp logs \u2014 camera frame to dashboard"],
            ["Safety Alert Accuracy", "\u226585%", "Reviewed alert samples over 30-day period"],
            ["Violation Detection", "\u226580%", "Spot checks vs. manual patrol data"],
            ["System Uptime", "\u226599.5%", "Server monitoring \u2014 30-day rolling window"],
            ["Revenue Accuracy", "\u226598%", "Reconciliation with manual collection records"],
        ],
        col_widths=[Inches(3.5), Inches(2.0), Inches(6.8)],
    )

    # ------------------------------------------------------------------
    # SLIDE 11 — Traffic Challenge
    # ------------------------------------------------------------------
    add_content_slide(prs, "The Challenge \u2014 Traffic Congestion", [
        "**Growing vehicle volume** \u2014 Bintulu's development is increasing road traffic year-on-year",
        "**Fixed-time signals** \u2014 no responsiveness to real-time traffic conditions",
        "**Global ATCS market** growing from $7B to $14B \u2014 proven demand for adaptive solutions",
        "**No real-time visibility** \u2014 traffic engineers lack live data for decision-making",
        "**Incident response delays** \u2014 no automated detection of stalled vehicles or accidents",
        "**Pedestrian safety** \u2014 fixed timing does not account for pedestrian flow patterns",
    ])

    # ------------------------------------------------------------------
    # SLIDE 12 — Traffic Architecture (diagram)
    # ------------------------------------------------------------------
    add_diagram_slide(prs, "AI Traffic Light \u2014 Solution Architecture",
                      _build_traffic_architecture)

    # ------------------------------------------------------------------
    # SLIDE 13 — Traffic Key Features
    # ------------------------------------------------------------------
    add_content_slide(prs, "AI Traffic Light \u2014 Key Features", [
        "**Vehicle Detection & Classification** \u2014 6 vehicle classes, 90%+ detection accuracy",
        "**Adaptive Signal Timing** \u2014 real-time phase adjustment based on queue length, <2s response",
        "**Multi-Lane Monitoring** \u2014 per-lane queue measurement, flow rate, turning counts",
        "**Incident Detection** \u2014 stalled vehicles, wrong-way driving, abnormal congestion alerts",
        "**Emergency Vehicle Priority** \u2014 signal preemption for ambulances, fire trucks, police",
        "**Central Dashboard** \u2014 all intersections on one screen, historical analytics, reporting",
    ])

    # ------------------------------------------------------------------
    # SLIDE 14 — Proven Results Worldwide
    # ------------------------------------------------------------------
    add_table_slide(
        prs,
        "AI Traffic Control \u2014 Proven Results Worldwide",
        ["City / Region", "Documented Result"],
        [
            ["Pittsburgh, PA (USA)", "25% reduction in travel time across corridors"],
            ["Tucson, AZ (USA)", "23% delay reduction, 1.25M+ driver hours saved annually"],
            ["London (UK)", "30% travel time reduction on managed corridors"],
            ["Las Vegas, NV (USA)", "17% crash reduction at AI-managed intersections"],
            ["Florida (USA, statewide)", "State DOT-approved AI traffic deployment across multiple districts"],
        ],
        col_widths=[Inches(4.0), Inches(8.3)],
    )

    # ------------------------------------------------------------------
    # SLIDE 15 — Traffic Cost & Timeline
    # ------------------------------------------------------------------
    add_content_slide(prs, "AI Traffic Light \u2014 Cost & Timeline", [
        "**$20,000 per intersection** \u2014 55% below the $45K global market median",
        "**Shared platform** from Smart Parking reduces incremental software investment",
        "**Standard industrial hardware** \u2014 no proprietary lock-in, same edge AI boxes",
        "**NTCIP protocol compatible** \u2014 integrates with existing signal controllers",
        "**Phase 2 deployment: months 8\u201314** \u2014 overlaps with parking full rollout",
        "**Scalable** \u2014 add intersections incrementally at $20K each, no platform changes",
        "**Failsafe design** \u2014 automatic revert to fixed timing if any AI component fails",
    ])

    # ------------------------------------------------------------------
    # SLIDE 16 — Shared Platform Architecture (diagram)
    # ------------------------------------------------------------------
    add_diagram_slide(prs, "Shared Platform Architecture",
                      _build_shared_platform)

    # ------------------------------------------------------------------
    # SLIDE 17 — Why VIETSOL + ESP
    # ------------------------------------------------------------------
    add_table_slide(
        prs,
        "Why VIETSOL + ESP",
        ["Factor", "VIETSOL + ESP", "Typical Competitor"],
        [
            ["AI Licensing", "Zero fees (open-source)", "Per-device annual fees"],
            ["Hardware", "Standard industrial processors", "Proprietary lock-in"],
            ["CCTV Reuse", "Any RTSP/ONVIF camera", "Requires proprietary cameras"],
            ["Parking Cost", "$450K for 7,000 bays", "$2\u2013$5M (sensor-based)"],
            ["Traffic Cost", "$20K/intersection", "$45K median"],
            ["Privacy", "Edge processing, metadata only", "Cloud video processing"],
            ["Local Presence", "ESP (Bintulu) + VIETSOL", "Foreign vendor"],
            ["Customisation", "Fully configurable & open", "Black box"],
        ],
        col_widths=[Inches(2.8), Inches(4.5), Inches(5.0)],
    )

    # ------------------------------------------------------------------
    # SLIDE 18 — Investment Summary
    # ------------------------------------------------------------------
    add_table_slide(
        prs,
        "Investment Summary",
        ["Use Case", "Investment", "Scope", "Timeline"],
        [
            ["Smart Parking (UC2)", "$450,000", "7,000 bays end-to-end", "Months 1\u201310"],
            ["AI Traffic Light (UC1)", "$20,000 / intersection", "Per intersection", "Months 8\u201314"],
            ["Shared Platform", "Included", "Backend + Dashboard + Mobile", "Months 3\u20134"],
        ],
        highlight_rows=[0],
        col_widths=[Inches(3.0), Inches(3.2), Inches(3.8), Inches(2.3)],
        font_size=Pt(15),
        header_font_size=Pt(16),
    )

    # ------------------------------------------------------------------
    # SLIDE 19 — Next Steps
    # ------------------------------------------------------------------
    add_content_slide(prs, "Next Steps", [
        "\u2776  **Site survey and infrastructure assessment** \u2014 camera mounting points, network, power",
        "\u2777  **Pilot zone selection** \u2014 500\u20131,000 bays for initial deployment and validation",
        "\u2778  **Detailed project plan and agreement** \u2014 scope, milestones, acceptance criteria",
        "\u2779  **Kickoff within 4 weeks** of agreement \u2014 hardware ordering, AI training begins",
        "",
        "We are ready to begin immediately.",
    ])

    # ------------------------------------------------------------------
    # SLIDE 20 — Thank You
    # ------------------------------------------------------------------
    _ = add_title_slide(
        prs,
        "Thank You",
        "VIETSOL + ESP",
        extra_line="Contact: [Name] | [Email] | [Phone]",
    )

    # ------------------------------------------------------------------
    # Save
    # ------------------------------------------------------------------
    output_path = Path(__file__).parent / "bintulu_smart_city_deck.pptx"
    prs.save(str(output_path))
    print(f"Presentation saved to {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
